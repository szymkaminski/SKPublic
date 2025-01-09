
import sys
import os
import docx
import fitz  # PyMuPDF
from pptx import Presentation
import pandas as pd
import json

# Add Azure OpenAI package
from openai import AzureOpenAI
# from dotenv import load_dotenv

# from azure.storage.blob import BlobServiceClient, BlobClient, ContainerClient


def main(): 
        
    try: 
    
        # Get configuration settings 
        # azure_oai_endpoint = os.getenv("AZURE_OAI_ENDPOINT")
        # azure_oai_key = os.getenv("AZURE_OAI_KEY")
        # azure_oai_deployment = os.getenv("AZURE_OAI_DEPLOYMENT")
        # Get configuration settings 
        azure_oai_endpoint = 'https://afpmarketrisk.openai.azure.com/'
        azure_oai_key = 'f2ad9b8845d24348bafd7d369c828dee'
        azure_oai_deployment = 'AFPChat4oCreative'
        
        # Azure Blob Storage configuration
        # blob_service_client = BlobServiceClient(account_url="https://afpmarketrisk.blob.core.windows.net", credential=DefaultAzureCredential())
        # container_client = blob_service_client.get_container_client("afpcreative")

        # Initialize the Azure OpenAI client...
        # Initialize the Azure OpenAI client
        client = AzureOpenAI(
                azure_endpoint = azure_oai_endpoint, 
                api_key=azure_oai_key,  
                api_version="2024-02-15-preview"
                )

        # Function to get data from Azure Blob Storage
        # def get_data_from_blob_storage():
        #     blob_list = container_client.list_blobs()
        #     data = ""
        #     for blob in blob_list:
        #         blob_client = container_client.get_blob_client(blob)
        #         blob_data = blob_client.download_blob().readall()
        #        data += blob_data.decode('utf-8') + "\n"
        #    return data

        # Fetch data from Azure Blob Storage
        # blob_data = get_data_from_blob_storage()

        # Function to read .docx file content
        def read_docx(file_path):
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])


        # Function to read PDF file content
        # This function will not read scanned content - another library is necessary for that
        def read_pdf(file_path):
            pdf_text = ""
            with fitz.open(file_path) as pdf_document:
                for page_num in range(pdf_document.page_count):
                    page = pdf_document.load_page(page_num)
                    pdf_text += page.get_text()
            return pdf_text
        
        # Function to read .pptx file content
        def read_pptx(file_path):
            ppt = Presentation(file_path)
            ppt_text = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        ppt_text += shape.text + "\n"
            return ppt_text
        
        # Function to read Excel file content
        def read_excel(file_path):
            excel_data = pd.read_excel(file_path)
            return excel_data.to_json()
        # Save the JSON data to a file
            json_filename = os.path.splitext(os.path.basename(file_path))[0] + '.json'
            json_output_path = os.path.join('<output_folder_path>', json_filename)
            with open(json_output_path, 'w') as json_file:
                json.dump(json.loads(json_data), json_file, indent=4)
            return json_data

        # Function to get data from specific files on local disk
        def get_data_from_local_files(file_paths):
            data = ""
            for file_path in file_paths:
                if file_path.endswith('.docx'):
                    data += read_docx(file_path) + "\n"
                elif file_path.endswith('.pdf'):
                    data += read_pdf(file_path) + "\n"
                elif file_path.endswith('.pptx'):
                    data += read_pptx(file_path) + "\n"
                elif file_path.endswith('.xlsx') or file_path.endswith('.xls'):
                    data += read_excel(file_path) + "\n"
                else:
                    try:
                        with open(file_path, 'r', encoding='utf-8') as file:
                            data += file.read() + "\n"
                    except UnicodeDecodeError:
                        with open(file_path, 'r', encoding='ISO-8859-1') as file:
                            data += file.read() + "\n"
            return data

        # List of specific files to use
        file_directory = r"C:\Users\skaminski\Documents\03. Projekty\Aramco\Python\Inputfiles"
#        file_directory = r"C:\Users\skaminski\Deloitte (O365D)\Projekty - General\Aramco - Market Risk - 2024\3. Working files\Notatki\Notatki przesłane do klienta"
        specific_files = [os.path.join(file_directory, f) for f in os.listdir(file_directory) if os.path.isfile(os.path.join(file_directory, f))]

        # Fetch data from specified files on local disk
        file_data = get_data_from_local_files(specific_files)
        
        Instruments_directory=r"C:\Users\skaminski\Documents\03. Projekty\Aramco\Python\xxx"
        instruments_data = Instruments_directory

        # add a sample conversation
        sample_conversation = """
            **User:** Please provide a template for the strategy for crude oil proces management

            **Assistant:** 1. Goal of strategy: reduce the volatility of EBITDA
            2. Detailed description of strategy: please provide a description based on the knowledge you acquired by reading
            the database \n\n{file_data}\n\n please inspire yourself by the description ("zasady zabezpieczania") 
            of the strategies of PKN Orlen, as well as other files uploaded into the database: \n\n{file_data}\n\n
            3. Hedging execution
            4. Hedging limits: at least a maximum and minimum volume limits should be defined, and potentially 
            an additional limit linked to the Company's risk appetite or target refining margin 
            5. Instruments used: please rely on the list of instruments provided in the file: \n\n{instruments_data}\n\n
            6. Hedging tenor: please provide hedging tenor. Please allow for the possibility of rollover, 
            as well as earlier hedging position liquidation in the case when the hedged item occurs earlier than expected
            7. Reporting 
            """
        # Create a system message
        system_message = f"""You are an expert in commodity and FX risk management. 
        You support Aramco Fuels Poland in designing its risk management strategy. You are using a formal, procedural language.
        Use the following data as your source of information:\n\n{file_data}\n\n
        In order to assess the specific situation of Aramco Fuels Poland, please rely on the report
        "Review of AFP’s financial risk database and determination of commodity and FX risk exposure"
        You should behave and respond in a manner consistent with the 
        following sample conversation:\n\n{sample_conversation}"""
        # Use the following data as your source of information:\n\n{blob_data}"""

        # Initialize messages array
        messages_array = [{"role": "system", "content": system_message}]
        
        # Create a Word document to save the conversation history
        word_file_path = 'conversation_history.docx'
        if not os.path.exists(word_file_path):
            doc = docx.Document()
            doc.save(word_file_path)
        else:
            doc = docx.Document(word_file_path)


        while True:
            # Get input text
            input_text = input("Enter the prompt (or type 'quit' to exit): ")
            if input_text.lower() == "quit":
                break
            if len(input_text) == 0:
                print("Please enter a prompt.")
                continue

            print("\nSending request for summary to Azure OpenAI endpoint...\n\n")
            
            # Send request to Azure OpenAI model
            messages_array.append({"role": "user", "content": input_text})

            response = client.chat.completions.create(
                model=azure_oai_deployment,
                temperature=0.5,            # Adjusts randomness of the output
                max_tokens=1200,             # Maximum number of tokens to generate
                top_p=0.5,                   # Controls diversity via nucleus sampling
#                frequency_penalty=0.0,       # Penalizes new tokens based on their frequency in the text so far
#                presence_penalty=0.0,        # Penalizes new tokens based on whether they appear in the text so far
                n=1,                         # Number of completions to generate for each prompt
                messages=messages_array
            )
            generated_text = response.choices[0].message.content
            # Add generated text to messages array
            messages_array.append({"role": "assistant", "content": generated_text})

            # Print generated text
            print("Summary: " + generated_text + "\n")
            
            # Save the input prompt and response to the Word file
            doc.add_paragraph(f"Prompt: {input_text}")
            doc.add_paragraph(f"Response: {generated_text}")
            doc.add_paragraph("\n" + "-"*50 + "\n")  # Add a separator
            doc.save(word_file_path)

    except Exception as ex:
        print(ex)

if __name__ == '__main__': 
    main()